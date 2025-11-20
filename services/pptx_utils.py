"""Helpers for PPTX post-processing such as hyperlink injection."""
from __future__ import annotations

from urllib.parse import quote_plus

from pptx import Presentation


def hyperlink_text_runs(
    pres: Presentation,
    needle: str,
    url: str,
    case_insensitive: bool = True,
    max_links: int = 3,
) -> int:
    """Attach a hyperlink to runs matching ``needle`` across the presentation."""

    if not needle or not url:
        return 0
    count = 0
    target = needle.lower() if case_insensitive else needle
    for slide in pres.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text_frame = shape.text_frame
            if not text_frame:
                continue
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text or ""
                    candidate = text.lower() if case_insensitive else text
                    match = candidate == target if case_insensitive else text == target
                    if not match:
                        continue
                    try:
                        run.hyperlink.address = url
                        count += 1
                        if count >= max_links:
                            return count
                    except Exception:
                        continue
    return count


def build_google_maps_url_from_address(address: str) -> str:
    query = quote_plus(address)
    return f"https://www.google.com/maps/search/?api=1&query={query}"


def apply_address_link(pres: Presentation, address_text: str, max_links: int = 3) -> int:
    """Apply a Google Maps hyperlink to occurrences of the provided address."""

    if not address_text:
        return 0
    try:
        url = build_google_maps_url_from_address(address_text)
        placed = hyperlink_text_runs(
            pres,
            address_text,
            url,
            case_insensitive=False,
            max_links=max_links,
        )
        if placed == 0:
            placed = hyperlink_text_runs(
                pres,
                address_text,
                url,
                case_insensitive=True,
                max_links=max_links,
            )
        return placed
    except Exception:
        return 0

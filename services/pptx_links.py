from urllib.parse import quote_plus
from pptx import Presentation


def add_hyperlink_to_text(prs: Presentation, exact_text: str, url: str) -> int:
    """
    Ajoute un lien cliquable sur TOUT RUN qui contient exactement `exact_text`.
    Retourne le nombre de runs modifiés.
    """
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip() == exact_text.strip():
                        h = run.hyperlink
                        h.address = url
                        count += 1
    return count

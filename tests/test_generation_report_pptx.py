import os
import sys
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation

ROOT_DIR = Path(__file__).resolve().parents[1]
if str(ROOT_DIR) not in sys.path:
    sys.path.insert(0, str(ROOT_DIR))

from app.services.pptx_fill import generate_book_pptx


def _create_temp_image() -> str:
    fd, path = tempfile.mkstemp(suffix=".png")
    os.close(fd)
    Image.new("RGB", (10, 10), color="red").save(path)
    return path


def test_generation_report_collects_missing_shapes(tmp_path):
    template_path = tmp_path / "template.pptx"
    output_path = tmp_path / "out.pptx"

    prs = Presentation()
    prs.slide_width = 9144000
    prs.slide_height = 6858000
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(template_path)

    image_path = _create_temp_image()
    report = generate_book_pptx(
        str(template_path),
        str(output_path),
        mapping={"[[ADRESSE]]": "Test"},
        image_by_shape={"SHAPE_MANQUANTE": image_path},
    )

    assert output_path.exists()
    assert "SHAPE_MANQUANTE" in report.missing_shapes
    assert report.ok is True

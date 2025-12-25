from pptx import Presentation
from pptx.util import Inches

from app.services.pptx_fill import replace_text_preserving_style
from app.services.token_utils import extract_pptx_tokens_from_presentation


def test_replace_tokens_in_table_and_group_shapes(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    table_shape = slide.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(2), Inches(1))
    table_shape.table.cell(0, 0).text = "[[TABLE_TOKEN]]"

    group_shape = slide.shapes.add_group_shape()
    group_shape.left = Inches(1)
    group_shape.top = Inches(2)
    group_shape.width = Inches(2)
    group_shape.height = Inches(1.5)
    textbox = group_shape.shapes.add_textbox(Inches(0.1), Inches(0.1), Inches(1.5), Inches(0.8))
    textbox.text = "[[GROUP_TOKEN]]"

    replace_text_preserving_style(
        slide.shapes,
        {"[[TABLE_TOKEN]]": "Table OK", "[[GROUP_TOKEN]]": "Group OK"},
    )

    out_path = tmp_path / "out.pptx"
    prs.save(out_path)
    after = Presentation(out_path)
    tokens_left = extract_pptx_tokens_from_presentation(after)

    assert "[[TABLE_TOKEN]]" not in tokens_left
    assert "[[GROUP_TOKEN]]" not in tokens_left

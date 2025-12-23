from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from app.services.template_validation import validate_pptx_template


def test_validate_pptx_template_missing_shape_and_unknown_token(tmp_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tx_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tx_box.text_frame.text = "[[UNKNOWN_TOKEN]]"
    path = tmp_path / "template.pptx"
    prs.save(path)

    result = validate_pptx_template(
        str(path),
        {"[[KNOWN_TOKEN]]"},
        required_shapes={"MAP_MASK"},
    )

    assert result.severity == "KO"
    assert "MAP_MASK" in result.missing_required_shapes
    assert "[[UNKNOWN_TOKEN]]" in result.unknown_tokens_in_template
    assert not result.ok

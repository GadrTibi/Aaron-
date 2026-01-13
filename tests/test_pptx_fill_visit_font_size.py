from pptx import Presentation
from pptx.util import Inches, Pt

from app.services.pptx_fill import replace_text_preserving_style


def test_visit_names_force_same_font_size() -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    box1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    p1 = box1.text_frame.paragraphs[0]
    run1a = p1.add_run()
    run1a.text = "[[VISITE_1_"
    run1a.font.size = Pt(18)
    run1b = p1.add_run()
    run1b.text = "NOM]]"
    run1b.font.size = Pt(12)

    box2 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    p2 = box2.text_frame.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "[[VISITE_2_NOM]]"
    run2.font.size = Pt(10)

    replace_text_preserving_style(
        slide.shapes,
        {
            "[[VISITE_1_NOM]]": "Louvre",
            "[[VISITE_2_NOM]]": "Montmartre",
        },
        force_font_size_tokens={"[[VISITE_1_NOM]]", "[[VISITE_2_NOM]]"},
        reference_token="[[VISITE_1_NOM]]",
    )

    sizes_1 = {run.font.size for run in p1.runs}
    sizes_2 = {run.font.size for run in p2.runs}
    assert sizes_1 == {Pt(18)}
    assert sizes_2 == {Pt(18)}

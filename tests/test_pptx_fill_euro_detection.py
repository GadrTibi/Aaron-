from pptx import Presentation
from pptx.util import Inches

from app.services.pptx_fill import detect_euro_in_template_context, replace_text_preserving_style


def test_detect_euro_in_template_context() -> None:
    text = "Prix par mois [[PRIX_MOIS]] €"
    assert detect_euro_in_template_context(text, "[[PRIX_MOIS]]") is True


def test_replace_strips_duplicate_euro() -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Prix par mois [[PRIX_MOIS]] €"

    replace_text_preserving_style(
        slide.shapes,
        {"[[PRIX_MOIS]]": "4150 €"},
    )

    assert "€€" not in box.text_frame.text
    assert box.text_frame.text == "Prix par mois 4150 €"

from pptx import Presentation
from pptx.util import Inches

from app.services.pptx_fill import replace_text_preserving_style


def _build_presentation_with_runs(run_texts):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    paragraph = textbox.text_frame.paragraphs[0]
    paragraph.clear()
    for text in run_texts:
        run = paragraph.add_run()
        run.text = text
    return prs, slide, textbox


def test_replace_tokens_preserves_split_runs_context(tmp_path):
    run_texts = [
        "Revenu mensuel brut généré : [[PRIX_NUI",
        "T]] la nuitée x [[JOURS_OCC_30]] jours/mois = [[REV_BRUT]]",
    ]
    mapping = {
        "[[PRIX_NUIT]]": "126€",
        "[[JOURS_OCC_30]]": "26",
        "[[REV_BRUT]]": "3.276 €",
    }
    prs, slide, textbox = _build_presentation_with_runs(run_texts)
    replace_text_preserving_style(slide.shapes, mapping)

    output_path = tmp_path / "out.pptx"
    prs.save(output_path)

    reloaded = Presentation(output_path)
    reloaded_textbox = reloaded.slides[0].shapes[0]
    assert (
        reloaded_textbox.text_frame.text
        == "Revenu mensuel brut généré : 126€ la nuitée x 26 jours/mois = 3.276 €"
    )


def test_replace_tokens_multiple_tokens_in_line(tmp_path):
    run_texts = [
        "*Frais de plateforme ([[PLATFORM_FEE_PCT]]% TTC = [[PLATFORM_FEE_EUR]]) + "
        "Commission MadeForYou ([[MFY_COMMISSION_PCT]]% TTC = [[MFY_COMMISSION_EUR]])",
    ]
    mapping = {
        "[[PLATFORM_FEE_PCT]]": "15",
        "[[PLATFORM_FEE_EUR]]": "492€",
        "[[MFY_COMMISSION_PCT]]": "15",
        "[[MFY_COMMISSION_EUR]]": "418€",
    }
    prs, slide, textbox = _build_presentation_with_runs(run_texts)
    replace_text_preserving_style(slide.shapes, mapping)

    output_path = tmp_path / "out.pptx"
    prs.save(output_path)

    reloaded = Presentation(output_path)
    reloaded_textbox = reloaded.slides[0].shapes[0]
    assert (
        reloaded_textbox.text_frame.text
        == "*Frais de plateforme (15% TTC = 492€) + Commission MadeForYou (15% TTC = 418€)"
    )

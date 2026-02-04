from __future__ import annotations

import os
from typing import Iterator

from pptx import Presentation

from app.services.token_utils import walk_pptx_shapes


def days_occupied_on_30(taux_occ_pct: float) -> float:
    return 30.0 * float(taux_occ_pct) / 100.0


def format_days(value: float, *, use_comma: bool = False) -> str:
    rounded = round(float(value), 1)
    if abs(rounded - round(rounded)) < 0.05:
        return str(int(round(rounded)))
    formatted = f"{rounded:.1f}"
    return formatted.replace(".", ",") if use_comma else formatted


def _iter_pptx_paragraph_texts(prs: Presentation) -> Iterator[str]:
    for slide in prs.slides:
        for shape in walk_pptx_shapes(slide.shapes):
            if hasattr(shape, "text_frame") and shape.text_frame:
                for para in shape.text_frame.paragraphs:
                    yield "".join(run.text or "" for run in para.runs)


def should_append_days_unit(
    template_path: str | None,
    *,
    token: str = "[[JOURS_OCC_30]]",
    word: str = "jour",
) -> tuple[bool, bool]:
    if not template_path or not os.path.exists(template_path):
        return True, False

    prs = Presentation(template_path)
    token_present = False
    context_has_word = False
    for text in _iter_pptx_paragraph_texts(prs):
        if token in text:
            token_present = True
            if word.lower() in text.lower():
                context_has_word = True
    return not context_has_word, token_present


def build_jours_occ_30_mapping(
    taux_occ_pct: float,
    *,
    days_value: float | None = None,
    include_unit: bool = True,
    use_comma: bool = False,
) -> dict[str, str]:
    days_num = days_occupied_on_30(taux_occ_pct) if days_value is None else float(days_value)
    days_str = format_days(days_num, use_comma=use_comma)
    value = f"{days_str} j" if include_unit else days_str
    return {"[[JOURS_OCC_30]]": value}

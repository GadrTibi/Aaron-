import os
from typing import Dict, Optional

import streamlit as st
from pptx import Presentation

from app.services.pptx_fill import replace_image_by_shape_name, replace_text_preserving_style


def generate_mandat_pptx(
    template_path: str,
    output_path: str,
    mapping: Dict[str, str],
    image_by_shape: Optional[Dict[str, str]] = None,
) -> None:
    mapping = dict(mapping) if mapping else {}
    mapping.setdefault(
        "[[MANDAT_DATE_SIGNATURE]]",
        st.session_state.get("mandat_date_signature_str", ""),
    )
    mapping.setdefault(
        "[[MANDAT_JOUR_SIGNATURE]]",
        st.session_state.get("mandat_jour_signature_str", ""),
    )

    prs = Presentation(template_path)
    for slide in prs.slides:
        replace_text_preserving_style(slide.shapes, mapping)

    if image_by_shape:
        for shape_name, img_path in image_by_shape.items():
            if not img_path:
                continue
            replace_image_by_shape_name(prs, shape_name, img_path)

    output_dir = os.path.dirname(output_path)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)
    prs.save(output_path)

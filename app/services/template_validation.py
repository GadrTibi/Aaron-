from __future__ import annotations

from dataclasses import dataclass, field
from typing import Callable, Dict, Iterable, Set

from docx import Document
from pptx import Presentation

from app.services.token_utils import (
    DOCX_TOKEN_PATTERN,
    PPTX_TOKEN_PATTERN,
    extract_docx_tokens_from_document,
    extract_pptx_tokens_from_presentation,
    extract_shape_names,
)


@dataclass
class ValidationResult:
    ok: bool
    severity: str
    missing_required_shapes: list[str] = field(default_factory=list)
    unknown_tokens_in_template: list[str] = field(default_factory=list)
    notes: list[str] = field(default_factory=list)


def _severity(missing_shapes: Iterable[str], unknown_tokens: Iterable[str]) -> str:
    if list(missing_shapes):
        return "KO"
    if list(unknown_tokens):
        return "WARN"
    return "OK"


def extract_docx_tokens(template_path: str) -> set[str]:
    doc = Document(template_path)
    return extract_docx_tokens_from_document(doc)


def extract_pptx_tokens(template_path: str) -> set[str]:
    prs = Presentation(template_path)
    return extract_pptx_tokens_from_presentation(prs)


def extract_pptx_shape_names(template_path: str) -> set[str]:
    prs = Presentation(template_path)
    names: Set[str] = set()
    for slide in prs.slides:
        names.update(extract_shape_names(slide.shapes))
    return names


def validate_docx_template(template_path: str, mapping_keys: set[str]) -> ValidationResult:
    tokens_in_template = extract_docx_tokens(template_path)
    unknown = sorted(tokens_in_template - mapping_keys)
    sev = _severity([], unknown)
    notes: list[str] = []
    if unknown:
        notes.append("Tokens inconnus détectés dans le template DOCX.")
    return ValidationResult(
        ok=sev != "KO",
        severity=sev,
        missing_required_shapes=[],
        unknown_tokens_in_template=unknown,
        notes=notes,
    )


def validate_pptx_template(
    template_path: str,
    mapping_keys: set[str],
    required_shapes: set[str],
    requirement_detectors: Dict[str, Callable[[Set[str]], bool]] | None = None,
) -> ValidationResult:
    tokens_in_template = extract_pptx_tokens(template_path)
    shapes_in_template = extract_pptx_shape_names(template_path)
    detectors = requirement_detectors or {}

    def _has_shape(req: str) -> bool:
        if req in shapes_in_template:
            return True
        detect_fn = detectors.get(req)
        if detect_fn:
            try:
                return bool(detect_fn(shapes_in_template))
            except Exception:
                return False
        return False

    unknown = sorted(tokens_in_template - mapping_keys)
    missing_shapes = sorted([req for req in required_shapes if not _has_shape(req)])
    sev = _severity(missing_shapes, unknown)
    notes: list[str] = []
    if missing_shapes:
        notes.append("Shapes requises absentes du template PPTX.")
    if unknown:
        notes.append("Tokens inconnus détectés dans le template PPTX.")
    return ValidationResult(
        ok=sev != "KO",
        severity=sev,
        missing_required_shapes=missing_shapes,
        unknown_tokens_in_template=unknown,
        notes=notes,
    )

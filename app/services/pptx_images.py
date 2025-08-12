import io, os
from typing import Optional, Tuple
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu
from PIL import Image, ImageDraw

# --- utils: normalisation de libellés ---
def _norm(s: str) -> str:
    return (s or "").replace("\u00A0", " ").strip().lower()

# --- marche tout: parcours récursif (groupes inclus) ---
def _walk_shapes_with_parent(slide):
    def rec(shapes, parent=None):
        for sh in shapes:
            yield sh, parent
            if hasattr(sh, "shapes"):
                yield from rec(sh.shapes, sh)
    yield from rec(slide.shapes, None)

def _get_alt_text(sh) -> Tuple[str, str]:
    el = getattr(sh, "element", None)
    if el is None: return "", ""
    # Alt-text dans Office: title (name) et descr (description)
    title = el.get("title") or ""
    descr = el.get("descr") or ""
    return title, descr

def find_shape_by_tag(prs: Presentation, tag: str) -> Tuple[Optional[object], Optional[object]]:
    """
    Cherche un shape dont le 'name' OU l'alt-text (title/descr) match 'tag' (insensible à la casse/espace).
    Parcourt toutes les slides et les enfants de groupes.
    Retourne (shape, slide) ou (None, None).
    """
    tagn = _norm(tag)
    for slide in prs.slides:
        for sh, _ in _walk_shapes_with_parent(slide):
            try:
                name_n = _norm(getattr(sh, "name", ""))
                t, d = _get_alt_text(sh)
                if name_n == tagn or _norm(t) == tagn or _norm(d) == tagn:
                    return sh, slide
            except Exception:
                pass
    return None, None

def _shape_bbox(sh) -> Tuple[Emu, Emu, Emu, Emu]:
    return sh.left, sh.top, sh.width, sh.height

def _to_circular_png(src_path: str, size_wh: Tuple[int,int]) -> bytes:
    # centre-crop en carré + redimension + masque circulaire alpha
    im = Image.open(src_path).convert("RGBA")
    w, h = im.size
    side = min(w, h)
    x0 = (w - side) // 2
    y0 = (h - side) // 2
    im = im.crop((x0, y0, x0 + side, y0 + side)).resize(size_wh, Image.LANCZOS)

    mask = Image.new("L", size_wh, 0)
    draw = ImageDraw.Draw(mask)
    draw.ellipse((0, 0, size_wh[0], size_wh[1]), fill=255)

    out = Image.new("RGBA", size_wh)
    out.paste(im, (0, 0), mask=mask)

    bio = io.BytesIO()
    out.save(bio, format="PNG")
    return bio.getvalue()

def inject_tagged_image(prs: Presentation, tag: str, image_path: str) -> bool:
    """
    1) Trouve le shape par nom/alt-text (tag).
    2) Si AUTO_SHAPE -> fill.user_picture(image_path) et return True.
    3) Sinon -> fallback: calcule bbox, fabrique PNG circulaire, add_picture à la même place.
    """
    sh, slide = find_shape_by_tag(prs, tag)
    if not sh or not slide:
        print(f"[WARN] Shape {tag} introuvable")
        return False

    left, top, width, height = _shape_bbox(sh)

    # 1) tentative masque natif
    try:
        if getattr(sh, "shape_type", None) == MSO_SHAPE_TYPE.AUTO_SHAPE and hasattr(sh, "fill"):
            sh.fill.user_picture(image_path)
            try: sh.fill.transparency = 0
            except Exception: pass
            print(f"[OK] Image injectée avec masque natif dans {tag}")
            return True
    except Exception:
        pass

    # 2) fallback garanti : PNG circulaire superposé
    try:
        # taille en pixels approx (EMU -> px ~ 96dpi)
        # On convertit l'EMU en pixels pour la fabrique du PNG
        # pptx utilise 914400 EMU par inch ; 96 px par inch
        def emu_to_px(e): return int((int(e) / 914400) * 96)
        wpx, hpx = max(2, emu_to_px(width)), max(2, emu_to_px(height))
        png_bytes = _to_circular_png(image_path, (wpx, hpx))
        # écrire PNG temporaire en mémoire
        bio = io.BytesIO(png_bytes)
        slide.shapes.add_picture(bio, left, top, width=width, height=height)
        print(f"[OK] Image injectée en fallback (PNG circulaire) dans {tag}")
        return True
    except Exception as e:
        print(f"[ERR] Fallback PNG circulaire a échoué pour {tag}: {e}")
        return False

# Compatibilité : ancien nom
inject_visit_image = inject_tagged_image

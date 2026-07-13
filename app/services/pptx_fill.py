import logging
import os
import re
from typing import Dict, Optional, List
from urllib.parse import quote_plus

from pptx import Presentation
from pptx.util import Inches
from PIL import Image

from app.services.pptx_images import inject_tagged_image
from app.services.pptx_links import add_hyperlink_to_text
from app.services.generation_report import GenerationReport
from app.services.token_utils import extract_pptx_tokens_from_presentation, walk_pptx_shapes

LOGGER = logging.getLogger(__name__)


def insert_plot_into_pptx(template_path: str, output_path: str, image_path: str, report: Optional[GenerationReport] = None, *, strict: bool = False) -> None:
    """Insert the histogram image into slide 6 using the dedicated mask shape."""

    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Fichier PPTX introuvable: {template_path}")
    if not os.path.exists(image_path):
        raise FileNotFoundError(
            f"Image d'histogramme introuvable: {image_path}. Générez le graphique avant l'export."
        )

    prs = Presentation(template_path)
    if len(prs.slides) <= 5:
        raise ValueError("La présentation ne contient pas de slide 6 pour y insérer l'histogramme.")

    slide = prs.slides[5]

    def _is_mask(name: Optional[str]) -> bool:
        if not name:
            return False
        norm = name.strip()
        if not norm:
            return False
        if norm.lower() == "estimation_histo_mask":
            return True
        return re.match(r"(?i).*histo.*-?mask$", norm) is not None

    target_shape = None
    for sh in walk_pptx_shapes(slide.shapes):
        if _is_mask(getattr(sh, "name", None)):
            target_shape = sh
            break

    if target_shape is None:
        message = "Shape mask introuvable en slide 6 (attendu: 'ESTIMATION_HISTO_MASK' ou variante '*histo*-mask')."
        if report is not None:
            report.add_missing_shapes(["ESTIMATION_HISTO_MASK"], blocking=strict)
            report.add_note(message)
            return
        raise ValueError(message)

    left, top, width, height = target_shape.left, target_shape.top, target_shape.width, target_shape.height

    output_dir = os.path.dirname(output_path)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)

    slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    LOGGER.info("Histogramme inséré dans la slide 6 (%s)", getattr(target_shape, "name", ""))
    prs.save(output_path)

def replace_tokens_in_paragraph(paragraph, mapping: Dict[str, str]) -> bool:
    runs = paragraph.runs
    if not runs:
        return False
    full = "".join(r.text for r in runs)
    if "[[" not in full:
        return False
    new = full
    for token, value in mapping.items():
        if token in new:
            new = new.replace(token, str(value))
    if new == full:
        return False
    runs[0].text = new
    for run in runs[1:]:
        run.text = ""
    return True

def replace_text_preserving_style(shapes, mapping: Dict[str, str]) -> None:
    for shape in walk_pptx_shapes(shapes):
        if hasattr(shape, "text_frame") and shape.text_frame:
            for para in shape.text_frame.paragraphs:
                replace_tokens_in_paragraph(para, mapping)

def insert_image(slide, image_path: str, left=Inches(1), top=Inches(3), width=Inches(8)) -> None:
    slide.shapes.add_picture(image_path, left, top, width=width)


def replace_image_by_shape_name(prs, shape_name: str, image_path: str, report: Optional[GenerationReport] = None, *, strict: bool = False) -> bool:
    """Remplace une image en la retrouvant par son nom."""
    try:
        ext = os.path.splitext(image_path)[1].lower()
        if ext not in (".jpg", ".jpeg", ".png"):
            tmp_path = os.path.splitext(image_path)[0] + ".png"
            with Image.open(image_path) as img:
                img.save(tmp_path, "PNG")
            image_path = tmp_path
    except Exception:
        pass

    for slide in prs.slides:
        for sh in walk_pptx_shapes(slide.shapes):
            try:
                if (sh.name or "").strip() == shape_name:
                    left, top, width, height = sh.left, sh.top, sh.width, sh.height
                    try:
                        sp = sh.element
                        sp.getparent().remove(sp)
                    except Exception:
                        pass
                    slide.shapes.add_picture(image_path, left, top, width=width, height=height)
                    if report is None:
                        print(f"[OK] Image remplacée dans {shape_name}")
                    return True
            except Exception:
                continue

    msg = f"Shape {shape_name} introuvable dans le PPTX."
    if report is not None:
        report.add_missing_shapes([shape_name], blocking=strict)
        report.add_note(msg)
    else:
        print(f"[WARN] {msg}")
    return False


def _collect_leftover_tokens(prs: Presentation) -> List[str]:
    return sorted(extract_pptx_tokens_from_presentation(prs))

def generate_estimation_pptx(
    template_path: str,
    output_path: str,
    mapping: Dict[str, str],
    chart_image: Optional[str] = None,
    image_by_shape: Optional[Dict[str, str]] = None,
    *,
    strict: bool = False,
) -> GenerationReport:
    report = GenerationReport()
    prs = Presentation(template_path)
    for slide in prs.slides:
        replace_text_preserving_style(slide.shapes, mapping)
    if image_by_shape:
        for shape_name, img_path in image_by_shape.items():
            if not img_path:
                continue
            # MAP_MASK doit toujours être injectée en plein rectangle sans masque
            if shape_name == "MAP_MASK":
                replace_image_by_shape_name(prs, shape_name, img_path, report, strict=strict)
            # VISITE_1_MASK et VISITE_2_MASK conservent l'injection avec masque circulaire
            elif shape_name in ("VISITE_1_MASK", "VISITE_2_MASK"):
                inject_tagged_image(prs, shape_name, img_path, report, strict=strict)
            else:
                replace_image_by_shape_name(prs, shape_name, img_path, report, strict=strict)
    output_dir = os.path.dirname(output_path)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)
    prs.save(output_path)

    leftovers = _collect_leftover_tokens(prs)
    if leftovers:
        report.add_missing_tokens(leftovers, blocking=strict)
        report.add_note("Des tokens sont restés dans le PPTX.")

    if chart_image:
        insert_plot_into_pptx(output_path, output_path, chart_image, report, strict=strict)
    return report


def generate_book_pptx(
    template_path: str,
    output_path: str,
    mapping: Dict[str, str],
    image_by_shape: Optional[Dict[str, str]] = None,
    *,
    strict: bool = False,
) -> GenerationReport:
    """Generate a Book PPTX from template and mapping/images."""
    report = GenerationReport()
    prs = Presentation(template_path)
    for slide in prs.slides:
        replace_text_preserving_style(slide.shapes, mapping)

    for token in ("[[ADRESSE]]", "[[BOOK_ADRESSE]]"):
        adresse = mapping.get(token, "").strip()
        if not adresse:
            continue
        maps_url = "https://www.google.com/maps/search/?api=1&query=" + quote_plus(adresse)
        add_hyperlink_to_text(prs, adresse, maps_url)

    if image_by_shape:
        for shape_name, img_path in image_by_shape.items():
            if not img_path:
                continue
            if shape_name in ("MAP_BOOK_MASK", "BOOK_MAP_MASK"):
                # Always replace the shape by the map image (full rectangle)
                replace_image_by_shape_name(prs, shape_name, img_path, report, strict=strict)
            elif shape_name in (
                "PORTE_ENTREE_MASK",
                "ENTREE_MASK",
                "APPARTEMENT_MASK",
                # Anciennes conventions
                "BOOK_ACCESS_PHOTO_PORTE",
                "BOOK_ACCESS_PHOTO_ENTREE",
                "BOOK_ACCESS_PHOTO_APPART",
            ):
                # Try native mask injection first, fallback to rectangle replace
                ok = inject_tagged_image(prs, shape_name, img_path, report, strict=strict)
                if not ok:
                    replace_image_by_shape_name(prs, shape_name, img_path, report, strict=strict)
            else:
                replace_image_by_shape_name(prs, shape_name, img_path, report, strict=strict)
    leftovers = _collect_leftover_tokens(prs)
    if leftovers:
        report.add_missing_tokens(leftovers, blocking=strict)
        report.add_note("Des tokens sont restés dans le PPTX.")
    prs.save(output_path)
    return report

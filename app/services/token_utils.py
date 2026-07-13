import re
from typing import Iterable, Iterator

from docx import Document
from pptx import Presentation


# Un token Word MFY est un PLACEHOLDER de fusion entre guillemets français. On code
# l'INTENTION, pas une forme figée (règle « conformité = intention, pas liste de
# tournures ») : les termes juridiques du corps du mandat sont écrits avec des
# guillemets TYPOGRAPHIQUES, donc un espace (normal ou insécable) COLLÉ au guillemet
# — « Mandant », « Bien », « Mandataire », « Notice d'Information… ». Un placeholder,
# lui, ne colle jamais d'espace au guillemet. Le motif exige donc un premier ET un
# dernier caractère NON-espace : il exclut les termes juridiques tout en détectant
# un placeholder même à espaces internes (ex. « Nom du gérant ») — contrairement à
# «\w+» qui rendrait un tel token INVISIBLE au garde-fou « document incomplet ».
DOCX_TOKEN_PATTERN = re.compile(r"«\S(?:[^»]*\S)?»")
PPTX_TOKEN_PATTERN = re.compile(r"\[\[[^\]]+\]\]")


def _collect_docx_paragraph_tokens(paragraph, pattern: re.Pattern[str]) -> set[str]:
    txt = "".join(run.text or "" for run in paragraph.runs)
    return set(pattern.findall(txt))


def extract_docx_tokens_from_document(doc: Document) -> set[str]:
    tokens: set[str] = set()

    def collect(paragraph) -> None:
        tokens.update(_collect_docx_paragraph_tokens(paragraph, DOCX_TOKEN_PATTERN))

    for p in doc.paragraphs:
        collect(p)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    collect(p)
    return tokens


def extract_docx_tokens(template_path: str) -> set[str]:
    """Retourne les tokens DOCX présents dans un fichier.

    Les tokens peuvent être fragmentés sur plusieurs *runs* Word ; la
    concaténation des runs de chaque paragraphe permet de les détecter
    correctement.
    """

    doc = Document(template_path)
    return extract_docx_tokens_from_document(doc)


def walk_pptx_shapes(shapes) -> Iterator[object]:
    """Yield all shapes recursively, diving into groups if present."""

    for sh in shapes:
        yield sh
        if hasattr(sh, "shapes"):
            yield from walk_pptx_shapes(sh.shapes)


def _collect_pptx_paragraph_tokens(paragraph, pattern: re.Pattern[str]) -> set[str]:
    txt = "".join(run.text or "" for run in paragraph.runs)
    return set(pattern.findall(txt))


def extract_pptx_tokens_from_presentation(prs: Presentation) -> set[str]:
    tokens: set[str] = set()
    for slide in prs.slides:
        for sh in walk_pptx_shapes(slide.shapes):
            if hasattr(sh, "text_frame") and sh.text_frame:
                for para in sh.text_frame.paragraphs:
                    tokens.update(_collect_pptx_paragraph_tokens(para, PPTX_TOKEN_PATTERN))
    return tokens


def extract_shape_names(shapes: Iterable[object]) -> set[str]:
    names: set[str] = set()
    for sh in walk_pptx_shapes(shapes):
        name = getattr(sh, "name", None)
        if name:
            norm = name.strip()
            if norm:
                names.add(norm)
    return names
